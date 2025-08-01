# /scripts/pdf_conversion.py (Corrected and with all comments restored)

# =============================================================================
# SECTION 1: DEPENDENCIES
# =============================================================================
import os
import sys
import json
import base64
import io
import httpx
import pandas as pd
import pdfplumber
from PIL import Image
import pytesseract
from pdf2image import convert_from_path
import camelot
from docx import Document
from pptx import Presentation
from openai import OpenAI
import traceback

# Import utilities
from logger import setup_logger
from ai_utils import get_ai_client, correct_text_with_ai, get_ai_quality_check

logger = setup_logger('pdf_conversion')

# =============================================================================
# SECTION 2: GLOBAL VARIABLES & INITIALIZATION
# =============================================================================
client = None
total_tokens_used_session = 0

def init_api_client():
    """Initializes the OpenAI client for AIMLAPI from environment variables."""
    global client
    api_key = os.environ.get('AIML_API_KEY')
    if not api_key:
        logger.error("Failed to initialize AIMLAPI client: AIML_API_KEY not found in environment.")
        client = None
        return None
    client = get_ai_client(api_key)
    logger.info("AIMLAPI client initialized for PDF conversion.")
    return client

# =============================================================================
# SECTION 3: HELPER FUNCTIONS
# =============================================================================

def image_to_base64(image: Image.Image) -> str:
    """Converts a PIL Image to a base64 encoded string."""
    buffered = io.BytesIO()
    image.save(buffered, format="PNG")
    return base64.b64encode(buffered.getvalue()).decode('utf-8')

def extract_text_from_page(pdf_path, page_num):
    """Extracts text from a single PDF page, using OCR as a fallback."""
    global total_tokens_used_session
    with pdfplumber.open(pdf_path) as pdf:
        page = pdf.pages[page_num - 1]
        text = page.extract_text()
        if text and text.strip():
            logger.info(f"Extracted selectable text from page {page_num}.")
            return text

        logger.warning(f"No selectable text on page {page_num}. Falling back to OCR.")
        try:
            images = convert_from_path(pdf_path, first_page=page_num, last_page=page_num, dpi=300)
            if not images: return ""
            
            ocr_text = pytesseract.image_to_string(images[0])
            if ocr_text and ocr_text.strip():
                logger.info(f"OCR successful for page {page_num}. Correcting with AI.")
                corrected_text, tokens_used = correct_text_with_ai(ocr_text)
                total_tokens_used_session += tokens_used
                return corrected_text
            return ""
        except Exception as e:
            logger.error(f"OCR process failed for page {page_num}: {traceback.format_exc()}")
            return ""

def extract_tables_with_vision(pdf_path, page_num):
    """Uses an AI vision model to extract tables from a PDF page image."""
    global total_tokens_used_session
    if not client: return []
    try:
        model_name = os.getenv('AIMLAPI_MODEL_VISION', "gpt-4o")
        images = convert_from_path(pdf_path, first_page=page_num, last_page=page_num, dpi=300)
        if not images: return []
        
        base64_image = image_to_base64(images[0])
        response = client.chat.completions.create(
            model=model_name,
            messages=[{"role": "user", "content": [{"type": "text", "text": "Extract all tables from this image into clean, pipe-delimited Markdown. If no tables are found, respond with 'No tables found'."}, {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}]}],
            max_tokens=2000,
        )
        total_tokens_used_session += response.usage.total_tokens
        logger.info(f"AI vision table extraction used {response.usage.total_tokens} tokens.")

        content = response.choices[0].message.content
        if "no tables found" in content.lower(): return []
        
        tables = []
        for table_md in [m.strip() for m in content.split("\n\n") if m.strip().startswith('|')]:
            try:
                lines = [l.strip() for l in table_md.split("\n") if l.strip()]
                if len(lines) < 2: continue
                header = [h.strip() for h in lines[0].split('|') if h.strip()]
                rows = [[td.strip() for td in line.split('|') if td.strip()] for line in lines[2:]]
                df = pd.DataFrame(rows, columns=header)
                tables.append(df)
            except Exception as e:
                logger.warning(f"Could not parse a markdown table from AI response: {e}")
        return tables
    except Exception as e:
        logger.error(f"AI Vision table extraction failed: {traceback.format_exc()}")
        return []

def extract_all_tables(pdf_path):
    """Extracts all tables from a PDF using multiple methods."""
    all_tables = []
    with pdfplumber.open(pdf_path) as pdf:
        for i in range(len(pdf.pages)):
            page_num = i + 1
            page_tables = []
            try:
                # Method 1: Camelot
                tables = camelot.read_pdf(pdf_path, pages=str(page_num), flavor='lattice')
                if tables.n > 0: page_tables.extend([tbl.df for tbl in tables])
                
                # Method 2: AI Vision as fallback
                if not page_tables:
                    logger.info(f"No tables found with rule-based methods on page {page_num}. Trying AI Vision.")
                    page_tables = extract_tables_with_vision(pdf_path, page_num)
                
                if page_tables:
                    all_tables.extend(page_tables)
            except Exception as e:
                logger.error(f"Table extraction failed on page {page_num}: {traceback.format_exc()}")
    return all_tables

# =============================================================================
# SECTION 4: MAIN CONVERSION FUNCTION
# =============================================================================

def convert_pdf(input_path, output_format, output_dir):
    """
    Main function to convert a PDF to a specified format.
    Returns the path to the primary output file and the tokens used.
    """
    global total_tokens_used_session
    total_tokens_used_session = 0 # Reset for each new file conversion
    
    base_name = os.path.splitext(os.path.basename(input_path))[0]
    output_path = os.path.join(output_dir, f"{base_name}.{output_format}")

    # --- Table-based Conversions ---
    if output_format in ['csv', 'xlsx']:
        tables = extract_all_tables(input_path)
        if not tables:
            raise ValueError(f"No tables could be extracted from {os.path.basename(input_path)}.")
        
        if output_format == 'xlsx':
            with pd.ExcelWriter(output_path, engine='xlsxwriter') as writer:
                for i, df in enumerate(tables):
                    df.to_excel(writer, sheet_name=f'Table_{i+1}', index=False)
        else: # For CSV, multiple files may be created
            for i, df in enumerate(tables):
                df.to_csv(os.path.join(output_dir, f"{base_name}_table_{i+1}.csv"), index=False)
    
    # --- Image-based Conversion ---
    elif output_format == 'png':
        images = convert_from_path(input_path, dpi=300)
        for i, image in enumerate(images):
            image.save(os.path.join(output_dir, f"{base_name}_page_{i+1}.png"), 'PNG')
    
    # --- Text-based Conversions ---
    else:
        full_text = ""
        with pdfplumber.open(input_path) as pdf:
            for i in range(len(pdf.pages)):
                full_text += extract_text_from_page(input_path, i + 1) + "\n\n"
        
        if output_format == 'docx':
            doc = Document()
            doc.add_paragraph(full_text)
            doc.save(output_path)
        elif output_format == 'pptx':
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            txBox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
            tf = txBox.text_frame
            tf.text = full_text
            tf.word_wrap = True
            prs.save(output_path)
        elif output_format == 'html':
            html_content = f"<html><head><title>{base_name}</title></head><body><pre>{full_text}</pre></body></html>"
            with open(output_path, 'w', encoding='utf-8') as f: 
                f.write(html_content)
        else: # txt, md
            with open(output_path, 'w', encoding='utf-8') as f: 
                f.write(full_text)

    return output_path, total_tokens_used_session

# =============================================================================
# SECTION 5: SCRIPT EXECUTION BLOCK
# =============================================================================
if __name__ == '__main__':
    try:
        output_dir_arg = sys.argv[1]
        output_format_arg = sys.argv[2]
        input_paths_arg = sys.argv[3:]
            
        init_api_client()
        
        total_tokens_all_files = 0
        for path in input_paths_arg:
            _, tokens_for_file = convert_pdf(path, output_format_arg, output_dir_arg)
            total_tokens_all_files += tokens_for_file

        print(json.dumps({"success": True, "tokens_used": total_tokens_all_files}))

    except ValueError as ve:
        # Handle specific known errors, like no tables being found
        if "No tables could be extracted" in str(ve):
            logger.error(f"No tables error: {traceback.format_exc()}")
            print(json.dumps({"error": "no_tables", "message": "No tables found in PDF. Try another file or different format.", "tokens_used": total_tokens_used_session}), file=sys.stderr)
        else:
            logger.error(f"Value error: {traceback.format_exc()}")
            print(json.dumps({"error": str(ve), "tokens_used": total_tokens_used_session}), file=sys.stderr)
        sys.exit(1)

    except Exception as e:
        # Handle all other unexpected errors
        logger.error(f"Main execution block error: {traceback.format_exc()}")
        print(json.dumps({"error": str(e), "tokens_used": total_tokens_used_session}), file=sys.stderr)
        sys.exit(1)